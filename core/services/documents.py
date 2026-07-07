"""Document text extraction (§5.1 — Premium, 3 generations/request).

Supported (per /help): pdf, docx, doc, xlsx, xls, csv, pptx, txt — up to 10 MB.
Each parser is imported lazily so a missing optional dependency degrades to a
clear error instead of breaking the bot."""
from __future__ import annotations

import csv
import io
import zipfile

MAX_FILE_BYTES = 10 * 1024 * 1024  # 10 MB
MAX_CHARS = 30_000  # cap extracted text fed to the model
# Zip-bomb guard: docx/xlsx/pptx are ZIP archives. Refuse if the declared
# uncompressed size is huge or the compression ratio is implausible.
MAX_DECOMPRESSED_BYTES = 200 * 1024 * 1024  # 200 MB total uncompressed
MAX_COMPRESSION_RATIO = 200

SUPPORTED_EXT = {"pdf", "docx", "doc", "xlsx", "xls", "csv", "pptx", "ppt", "txt"}


class UnsupportedDocument(Exception):
    pass


def ext_of(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def _guard_zip(data: bytes) -> io.BytesIO:
    """Validate a ZIP-based Office file against zip-bomb expansion before any
    parser touches it. Returns a fresh BytesIO ready for the parser."""
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            total = sum(info.file_size for info in zf.infolist())
    except zipfile.BadZipFile as exc:
        raise UnsupportedDocument(f"corrupt archive: {exc}") from exc
    if total > MAX_DECOMPRESSED_BYTES:
        raise UnsupportedDocument("document expands too large (possible zip bomb)")
    if data and total / len(data) > MAX_COMPRESSION_RATIO:
        raise UnsupportedDocument("compression ratio too high (possible zip bomb)")
    return io.BytesIO(data)


def _pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _docx(data: bytes) -> str:
    from docx import Document

    doc = Document(_guard_zip(data))
    return "\n".join(p.text for p in doc.paragraphs)


def _xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(_guard_zip(data), read_only=True, data_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"# {ws.title}")
        for row in ws.iter_rows(values_only=True):
            out.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(out)


def _pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(_guard_zip(data))
    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return "\n".join(out)


def _csv(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    rows = csv.reader(io.StringIO(text))
    return "\n".join("\t".join(r) for r in rows)


def _txt(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


_EXTRACTORS = {
    "pdf": _pdf,
    "docx": _docx,
    "doc": _docx,
    "xlsx": _xlsx,
    "xls": _xlsx,
    "pptx": _pptx,
    "ppt": _pptx,
    "csv": _csv,
    "txt": _txt,
}


def extract_text(filename: str, data: bytes) -> str:
    ext = ext_of(filename)
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise UnsupportedDocument(ext)
    try:
        text = extractor(data)
    except ImportError as exc:  # parser lib not installed
        raise UnsupportedDocument(f"{ext}: parser unavailable ({exc.name})") from exc
    return text.strip()[:MAX_CHARS]
